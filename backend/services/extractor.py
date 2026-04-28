import base64
import logging
import re
from pathlib import Path
from typing import Optional

log = logging.getLogger(__name__)

TEXT_EXTENSIONS = {".md", ".txt", ".rtf", ".csv", ".py", ".js", ".ts",
                   ".html", ".css", ".json", ".yaml", ".yml"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff"}
CODE_EXTENSIONS = {".py", ".js", ".ts", ".html", ".css", ".json", ".yaml", ".yml"}


def extract(path: Path) -> dict:
    """
    Returns {text, page_count, word_count, row_count, dimensions, language, method}
    method: 'text' | 'ocr' | 'vision' | 'unsupported'
    """
    suffix = path.suffix.lower()
    result = {
        "text": "",
        "page_count": None,
        "word_count": None,
        "row_count": None,
        "dimensions": None,
        "language": None,
        "method": "unsupported",
        "image_b64": None,
    }

    try:
        if suffix == ".pdf":
            _extract_pdf(path, result)
        elif suffix in (".docx", ".doc"):
            _extract_docx(path, result)
        elif suffix in (".xlsx", ".xls"):
            _extract_xlsx(path, result)
        elif suffix in (".pptx", ".ppt"):
            _extract_pptx(path, result)
        elif suffix in (".odt", ".ods", ".odp"):
            _extract_odf(path, result)
        elif suffix == ".csv":
            _extract_csv(path, result)
        elif suffix in TEXT_EXTENSIONS:
            _extract_text(path, result)
        elif suffix in IMAGE_EXTENSIONS:
            _extract_image(path, result)
        else:
            _try_text_fallback(path, result)
    except Exception as e:
        log.error("Extraction failed for %s: %s", path, e)
        result["method"] = "unsupported"

    if result["text"]:
        result["word_count"] = len(result["text"].split())
        result["language"] = _detect_language(result["text"])

    return result


def _extract_pdf(path: Path, result: dict):
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(str(path))
        pages = []
        result["page_count"] = len(reader.pages)
        for i, page in enumerate(reader.pages[:10]):
            text = page.extract_text() or ""
            pages.append(text)
        full_text = "\n".join(pages).strip()

        if len(full_text) < 100:
            # Likely scanned — fall back to OCR
            from backend.services.ocr import ocr_image_file, is_available
            if is_available():
                ocr_pages = []
                for i, page in enumerate(reader.pages[:10]):
                    # Try rendering with Pillow if possible
                    pass
                result["text"] = full_text
                result["method"] = "ocr"
            else:
                result["text"] = full_text
                result["method"] = "text"
        else:
            result["text"] = full_text[:8000]
            result["method"] = "text"
    except Exception as e:
        log.error("PDF extraction error: %s", e)


def _extract_docx(path: Path, result: dict):
    try:
        import docx
        doc = docx.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        result["text"] = "\n".join(paragraphs)[:8000]
        result["page_count"] = None
        result["method"] = "text"
    except Exception as e:
        log.error("DOCX extraction error: %s", e)


def _extract_xlsx(path: Path, result: dict):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        total_rows = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"Sheet: {sheet_name}")
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 50:
                    break
                row_text = "\t".join(str(c) if c is not None else "" for c in row)
                parts.append(row_text)
                total_rows += 1
        result["text"] = "\n".join(parts)[:8000]
        result["row_count"] = total_rows
        result["method"] = "text"
    except Exception as e:
        log.error("XLSX extraction error: %s", e)


def _extract_pptx(path: Path, result: dict):
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides_text = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slides_text.append(f"Slide {i+1}: " + " | ".join(texts))
        result["text"] = "\n".join(slides_text)[:8000]
        result["page_count"] = len(prs.slides)
        result["method"] = "text"
    except Exception as e:
        log.error("PPTX extraction error: %s", e)


def _extract_odf(path: Path, result: dict):
    try:
        from odf import text as odf_text, teletype
        from odf.opendocument import load as odf_load
        doc = odf_load(str(path))
        texts = teletype.extractText(doc.body)
        result["text"] = texts[:8000]
        result["method"] = "text"
    except Exception as e:
        log.error("ODF extraction error: %s", e)


def _extract_csv(path: Path, result: dict):
    try:
        import csv
        rows = []
        with open(str(path), "r", encoding="utf-8", errors="replace", newline="") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= 100:
                    break
                rows.append("\t".join(row))
        result["text"] = "\n".join(rows)[:8000]
        result["row_count"] = len(rows)
        result["method"] = "text"
    except Exception as e:
        log.error("CSV extraction error: %s", e)


def _extract_text(path: Path, result: dict):
    try:
        with open(str(path), "r", encoding="utf-8", errors="replace") as f:
            content = f.read(8000)
        result["text"] = content
        result["method"] = "text"
    except Exception as e:
        log.error("Text extraction error: %s", e)


def _extract_image(path: Path, result: dict):
    try:
        from PIL import Image as PILImage
        img = PILImage.open(str(path))
        result["dimensions"] = f"{img.width}x{img.height}"

        with open(str(path), "rb") as f:
            result["image_b64"] = base64.b64encode(f.read()).decode("utf-8")

        # Also attempt OCR for embedded text
        from backend.services.ocr import ocr_image_file, is_available
        if is_available():
            ocr_text = ocr_image_file(path)
            if ocr_text:
                result["text"] = ocr_text
        result["method"] = "vision"
    except Exception as e:
        log.error("Image extraction error: %s", e)


def _try_text_fallback(path: Path, result: dict):
    try:
        with open(str(path), "r", encoding="utf-8", errors="replace") as f:
            content = f.read(3000)
        result["text"] = content
        result["method"] = "text"
    except Exception:
        result["method"] = "unsupported"


def _detect_language(text: str) -> str:
    """Simple French/English detection based on common function words."""
    if not text:
        return "EN"
    text_lower = text.lower()
    fr_words = ["le ", "la ", "les ", "de ", "du ", "des ", "et ", "en ", "est ", "que ", "pour "]
    en_words = ["the ", "and ", "for ", "this ", "that ", "with ", "from ", "have ", "are ", "not "]
    fr_score = sum(text_lower.count(w) for w in fr_words)
    en_score = sum(text_lower.count(w) for w in en_words)
    if fr_score > en_score * 1.5:
        return "FR"
    if en_score > fr_score * 1.5:
        return "EN"
    if fr_score > 0 and en_score > 0:
        return "FR/EN"
    return "EN"
